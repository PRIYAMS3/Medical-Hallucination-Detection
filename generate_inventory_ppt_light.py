from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def p(points: float):
    return Inches(points / 72.0)


COLORS = {
    "bg": RGBColor(248, 250, 252),
    "panel": RGBColor(255, 255, 255),
    "panel_alt": RGBColor(241, 245, 249),
    "title": RGBColor(15, 23, 42),
    "text": RGBColor(51, 65, 85),
    "muted": RGBColor(100, 116, 139),
    "line": RGBColor(203, 213, 225),
    "accent": RGBColor(37, 99, 235),
}


def set_font(run, size, color, bold=False, name="Calibri"):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name


def add_background(slide, sw, sh):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["bg"]
    bg.line.fill.background()

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, p(6))
    top.fill.solid()
    top.fill.fore_color.rgb = COLORS["accent"]
    top.line.fill.background()


def add_textbox(slide, text, x, y, w, h, size=16, color=None, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(p(x), p(y), p(w), p(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    set_font(run, size, color or COLORS["text"], bold=bold)
    return box


def add_bullets(slide, items, x, y, w, h, size=18):
    box = slide.shapes.add_textbox(p(x), p(y), p(w), p(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, text in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = text
        para.level = 0
        para.space_after = Pt(8)
        para.bullet = True
        para.alignment = PP_ALIGN.LEFT
        for run in para.runs:
            set_font(run, size, COLORS["text"], bold=False)
    return box


def add_card(slide, x, y, w, h, title):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p(x), p(y), p(w), p(h))
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["panel"]
    card.line.color.rgb = COLORS["line"]
    card.adjustments[0] = 0.02
    add_textbox(slide, title, x + 16, y + 12, w - 26, 24, size=14, color=COLORS["muted"], bold=True)
    return card


def add_placeholder(slide, text, x, y, w=430, h=48):
    ph = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p(x), p(y), p(w), p(h))
    ph.fill.solid()
    ph.fill.fore_color.rgb = COLORS["panel_alt"]
    ph.line.color.rgb = COLORS["line"]
    ph.adjustments[0] = 0.03
    add_textbox(slide, f"[{text}]", x + 14, y + 14, w - 20, h - 10, size=14, color=COLORS["muted"], bold=True)


def add_title(slide, title, subtitle=""):
    add_textbox(slide, title, 44, 28, 860, 46, size=34, color=COLORS["title"], bold=True)
    if subtitle:
        add_textbox(slide, subtitle, 44, 76, 860, 26, size=16, color=COLORS["muted"])


def add_speaking_hint_box(slide, hint):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p(44), p(466), p(872), p(46))
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS["panel_alt"]
    box.line.color.rgb = COLORS["line"]
    box.adjustments[0] = 0.03
    add_textbox(slide, f"Presenter Hint: {hint}", 56, 480, 844, 24, size=13, color=COLORS["muted"])


def add_notes(slide, hint):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = f"Speaking Hint: {hint}"


def add_two_col_layout(slide, heading, bullets, placeholder_text=None):
    add_title(slide, heading)
    add_card(slide, 44, 116, 560, 334, "Key Points")
    add_bullets(slide, bullets, 62, 150, 524, 260, size=17)
    if placeholder_text:
        add_card(slide, 626, 116, 290, 334, "Placeholder")
        add_placeholder(slide, placeholder_text, 646, 250, 250, 70)


def build_presentation(output: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    sw = prs.slide_width
    sh = prs.slide_height

    # 1
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title(s, "AI Smart Inventory Management System", "Demand Forecasting using Machine Learning")
    add_card(s, 44, 146, 872, 258, "Project Introduction")
    add_bullets(
        s,
        [
            "Team Members: [Add 4 Names]",
            "Course / Institution: [Add Details]",
            "Final-year project integrating Machine Learning and MLOps practices",
        ],
        66,
        182,
        836,
        190,
        size=20,
    )
    hint = "Briefly introduce the project, team, and why this topic matters."
    add_speaking_hint_box(s, hint)
    add_notes(s, hint)

    # 2
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title(s, "Problem Statement")
    add_card(s, 44, 116, 872, 334, "Business Challenge")
    add_bullets(
        s,
        [
            "Retail businesses struggle with accurate demand prediction across stores and product categories.",
            "Overstocking increases storage and operational costs, reducing overall profitability.",
            "Understocking leads to lost revenue and poor customer experience due to stockouts.",
            "Traditional forecasting methods fail to capture dynamic demand patterns.",
            "Summary: Accurate demand forecasting is critical for efficient inventory management.",
        ],
        62,
        150,
        844,
        290,
        size=16,
    )
    hint = "Explain why demand forecasting is crucial in real retail operations and decision-making."
    add_speaking_hint_box(s, hint)
    add_notes(s, hint)

    # 3
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title(s, "Solution Overview")
    add_card(s, 44, 116, 872, 334, "Approach")
    add_bullets(
        s,
        [
            "Developed an ML-based demand forecasting system for inventory planning.",
            "Uses historical sales data and engineered features to learn demand behavior.",
            "Provides real-time predictions through API and a business-friendly dashboard.",
            "Designed as an end-to-end automated pipeline from training to deployment.",
        ],
        62,
        150,
        844,
        280,
        size=17,
    )
    hint = "Walk through the full solution in sequence: data, model, API, and dashboard."
    add_speaking_hint_box(s, hint)
    add_notes(s, hint)

    # 4
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_two_col_layout(
        s,
        "Dataset",
        [
            "Time-series retail dataset with daily observations.",
            "Key features include Store ID, Product family, Date, Sales, and Promotion.",
            "Each row represents sales activity for one store-product combination on a specific date.",
            "Dataset link is provided externally for reproducibility.",
        ],
        "Add dataset sample screenshot",
    )
    hint = "Clarify what one row means and how the dataset supports forecasting."
    add_speaking_hint_box(s, hint)
    add_notes(s, hint)

    # 5
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_two_col_layout(
        s,
        "Data Preprocessing",
        [
            "Converted date column into datetime format for time-aware processing.",
            "Sorted records chronologically (store -> family -> date) to preserve temporal order.",
            "Handled missing values: sales filled with 0 and promotion filled with 0.",
            "Encoded categorical variables for model compatibility.",
            "Key point: identical preprocessing logic is reused in training and inference.",
        ],
        "Add preprocess.py code snippet",
    )
    hint = "Emphasize consistency between training and prediction preprocessing."
    add_speaking_hint_box(s, hint)
    add_notes(s, hint)

    # 6
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_two_col_layout(
        s,
        "Feature Engineering (Core)",
        [
            "Time-based features: day_of_week, month, week, weekend.",
            "Lag features: lag_1, lag_7, lag_14, lag_28 to capture recent demand behavior.",
            "Rolling statistics: rolling_mean_7 and rolling_std_7 for trend and volatility.",
            "Insight: these features capture seasonality, historical momentum, and local fluctuations.",
        ],
        "Add features.py code snippet",
    )
    hint = "Explain why lag features are powerful: they let the model learn demand memory."
    add_speaking_hint_box(s, hint)
    add_notes(s, hint)

    # 7
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_two_col_layout(
        s,
        "Data Leakage Prevention",
        [
            "All predictor features are computed from past data only using shift(1).",
            "Target is defined as a future value using shift(-7) for 7-day-ahead forecasting.",
            "Strict chronological ordering is maintained throughout the pipeline.",
            "Importance: prevents unrealistic performance and ensures reliable evaluation.",
        ],
        "Add leakage prevention code snippet",
    )
    hint = "Mention this is a common pitfall and highlight that the project explicitly avoided it."
    add_speaking_hint_box(s, hint)
    add_notes(s, hint)

    # 8
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_two_col_layout(
        s,
        "Model Training",
        [
            "Trained two models: Random Forest and XGBoost.",
            "Applied an 80/20 train-test split to evaluate generalization.",
            "Performed hyperparameter tuning to improve model quality.",
            "Goal: choose the best model based on reliable evaluation metrics.",
        ],
        "Add training code snippet",
    )
    hint = "Briefly mention XGBoost is strong for tabular data due to boosted trees."
    add_speaking_hint_box(s, hint)
    add_notes(s, hint)

    # 9
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title(s, "Model Evaluation")
    add_card(s, 44, 116, 420, 334, "Evaluation Summary")
    add_bullets(
        s,
        [
            "Primary metric: RMSE (Root Mean Squared Error).",
            "Lower RMSE indicates better prediction accuracy.",
            "Conclusion: XGBoost was selected as the final model.",
        ],
        62,
        150,
        390,
        180,
        size=17,
    )
    table = s.shapes.add_table(3, 2, p(492), p(152), p(390), p(170)).table
    table.columns[0].width = p(220)
    table.columns[1].width = p(170)
    table.cell(0, 0).text = "Model"
    table.cell(0, 1).text = "Performance"
    table.cell(1, 0).text = "Random Forest"
    table.cell(1, 1).text = "Higher RMSE"
    table.cell(2, 0).text = "XGBoost"
    table.cell(2, 1).text = "Lower RMSE"
    for r in range(3):
        for c in range(2):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["panel_alt"] if r == 0 else COLORS["panel"]
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            for run in cell.text_frame.paragraphs[0].runs:
                set_font(run, 14, COLORS["title"] if r == 0 else COLORS["text"], bold=(r == 0))
    add_placeholder(s, "Add RMSE table/chart", 492, 350, 390, 60)
    hint = "Explain RMSE in simple terms: average prediction error magnitude in sales units."
    add_speaking_hint_box(s, hint)
    add_notes(s, hint)

    # 10
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_two_col_layout(
        s,
        "Training Pipeline",
        [
            "Flow: Data -> Preprocess -> Feature Engineering -> Train -> Save Model.",
            "Implemented inside pipeline.py for automated execution.",
            "Enables reproducible training with one command.",
        ],
        "Add pipeline code snippet",
    )
    hint = "State that one script runs the complete ML training workflow."
    add_speaking_hint_box(s, hint)
    add_notes(s, hint)

    # 11
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_card(s, 44, 116, 872, 334, "Inference Pipeline")
    add_title(s, "Inference Pipeline")
    add_bullets(
        s,
        [
            "Flow: Input -> Preprocess -> Feature Engineering -> Prediction.",
            "Inference reuses the same logic and transformations used during training.",
            "Ensures consistency, reliability, and reduced production errors.",
        ],
        62,
        150,
        844,
        180,
        size=17,
    )
    add_placeholder(s, "Add inference flow diagram", 62, 336, 844, 80)
    hint = "Highlight there is no mismatch between training and prediction pipelines."
    add_speaking_hint_box(s, hint)
    add_notes(s, hint)

    # 12
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_two_col_layout(
        s,
        "FastAPI Deployment",
        [
            "Model deployed as a REST API using FastAPI.",
            "Endpoint: POST /predict.",
            "Input format: JSON payload with required features.",
            "Output: predicted demand value for decision support.",
        ],
        "Add app.py endpoint code",
    )
    add_placeholder(s, "Add /docs screenshot", 646, 332, 250, 70)
    hint = "Explain that API deployment enables easy integration with external systems."
    add_speaking_hint_box(s, hint)
    add_notes(s, hint)

    # 13
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_two_col_layout(
        s,
        "Dashboard (Streamlit)",
        [
            "Interactive UI allows users to request predictions without coding.",
            "Users can select store and product family from input controls.",
            "Lag features are auto-generated in the backend pipeline.",
            "Charts show past demand against predicted demand for quick interpretation.",
        ],
        "Add Streamlit screenshot",
    )
    hint = "Position this as the business-facing interface for non-technical users."
    add_speaking_hint_box(s, hint)
    add_notes(s, hint)

    # 14
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_two_col_layout(
        s,
        "MLOps Components",
        [
            "Docker: containerized environment for consistent deployment.",
            "CI: automated checks and workflows using GitHub Actions.",
            "Logging: prediction and runtime logs stored in logs/app.log.",
            "Model versioning: tracked model artifact (model_v1.pkl).",
        ],
        "Add Docker + GitHub Actions screenshot",
    )
    hint = "Conclude this slide by saying these practices make the system production-ready."
    add_speaking_hint_box(s, hint)
    add_notes(s, hint)

    # 15
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_two_col_layout(
        s,
        "System Architecture",
        [
            "src: core ML logic and reusable components.",
            "pipeline: automation scripts for training workflow.",
            "app: FastAPI service for prediction endpoint.",
            "models: saved trained model artifacts.",
            "logs: runtime and prediction logging outputs.",
        ],
        "Add architecture diagram",
    )
    hint = "Explain each module and how data flows across components."
    add_speaking_hint_box(s, hint)
    add_notes(s, hint)

    # 16
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title(s, "Live Demo")
    add_card(s, 44, 116, 872, 334, "Demo Plan")
    add_bullets(
        s,
        [
            "Show FastAPI interactive docs at /docs.",
            "Run a sample prediction request and explain the JSON response.",
            "Open Streamlit dashboard and demonstrate end-user interaction.",
        ],
        62,
        164,
        844,
        200,
        size=18,
    )
    add_placeholder(s, "Add demo sequence screenshot(s)", 62, 356, 844, 62)
    hint = "Transition smoothly from slides to live system by stating the demo sequence."
    add_speaking_hint_box(s, hint)
    add_notes(s, hint)

    # 17
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title(s, "Conclusion")
    add_card(s, 44, 116, 872, 334, "Project Outcomes")
    add_bullets(
        s,
        [
            "Built an end-to-end ML system for inventory demand forecasting.",
            "Achieved accurate demand prediction with robust feature engineering and model selection.",
            "Designed the solution for scalability and practical real-world use.",
        ],
        62,
        164,
        844,
        210,
        size=18,
    )
    hint = "Summarize impact: technical completeness plus business value."
    add_speaking_hint_box(s, hint)
    add_notes(s, hint)

    # 18
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title(s, "Future Work")
    add_card(s, 44, 116, 872, 334, "Next Steps")
    add_bullets(
        s,
        [
            "Incorporate external signals such as weather and holidays.",
            "Explore deep learning models like LSTM for sequence forecasting.",
            "Enable real-time integration with enterprise inventory systems.",
        ],
        62,
        164,
        844,
        210,
        size=18,
    )
    hint = "Present future work as a roadmap from strong baseline to advanced deployment."
    add_speaking_hint_box(s, hint)
    add_notes(s, hint)

    # 19
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title(s, "Thank You")
    add_card(s, 180, 170, 600, 200, "Closing")
    add_textbox(s, "Thank you\nQuestions", 200, 220, 560, 110, size=34, color=COLORS["title"], bold=True, align=PP_ALIGN.CENTER)
    hint = "Invite questions and discussion."
    add_speaking_hint_box(s, hint)
    add_notes(s, hint)

    prs.save(str(output))


if __name__ == "__main__":
    out = Path.cwd() / "AI_Smart_Inventory_Management_System_Light.pptx"
    build_presentation(out)
    print(f"Created: {out}")
