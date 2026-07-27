from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def p(points: float):
    return Inches(points / 72.0)


COLORS = {
    "bg": RGBColor(11, 15, 23),
    "panel": RGBColor(18, 26, 39),
    "panel_soft": RGBColor(23, 34, 49),
    "white": RGBColor(244, 248, 255),
    "muted": RGBColor(143, 163, 191),
    "line": RGBColor(38, 50, 68),
    "accent": RGBColor(34, 211, 238),
    "accent2": RGBColor(103, 232, 167),
    "warn": RGBColor(250, 204, 21),
    "bad": RGBColor(248, 113, 113),
}


def set_font(run, size, color, bold=False, name="Segoe UI"):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name


def add_background(slide, sw, sh):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["bg"]
    bg.line.fill.background()

    orb1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, p(-120), p(-120), p(320), p(320))
    orb1.fill.solid()
    orb1.fill.fore_color.rgb = COLORS["accent"]
    orb1.fill.transparency = 0.90
    orb1.line.fill.background()

    orb2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, sw - p(220), sh - p(220), p(280), p(280))
    orb2.fill.solid()
    orb2.fill.fore_color.rgb = COLORS["accent2"]
    orb2.fill.transparency = 0.91
    orb2.line.fill.background()

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, p(5))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS["accent"]
    top_bar.line.fill.background()


def add_textbox(slide, text, x, y, w, h, size=14, color=None, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(p(x), p(y), p(w), p(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    set_font(run, size, color or COLORS["white"], bold=bold)
    return box


def add_card(slide, x, y, w, h, title, body):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p(x), p(y), p(w), p(h))
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["panel"]
    card.line.color.rgb = COLORS["line"]
    card.adjustments[0] = 0.08
    add_textbox(slide, title, x + 14, y + 10, w - 24, 26, size=16, bold=True, color=COLORS["white"])
    add_textbox(slide, body, x + 14, y + 42, w - 24, h - 50, size=12, color=COLORS["muted"])
    return card


def add_chip(slide, text, x, y, w=390, h=34, fill_key="panel_soft"):
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p(x), p(y), p(w), p(h))
    chip.fill.solid()
    chip.fill.fore_color.rgb = COLORS[fill_key]
    chip.line.color.rgb = COLORS["line"]
    chip.adjustments[0] = 0.2
    add_textbox(slide, text, x + 12, y + 7, w - 20, 22, size=13, bold=True, color=COLORS["white"])
    return chip


def add_title_block(slide, title, subtitle=""):
    add_textbox(slide, title, 44, 32, 870, 76, size=36, bold=True, color=COLORS["white"])
    if subtitle:
        add_textbox(slide, subtitle, 44, 90, 860, 34, size=16, color=COLORS["muted"])


def add_flow_step(slide, x, y, w, h, label, accent=False):
    step = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p(x), p(y), p(w), p(h))
    step.fill.solid()
    step.fill.fore_color.rgb = COLORS["panel_soft"] if accent else COLORS["panel"]
    step.line.color.rgb = COLORS["line"]
    step.adjustments[0] = 0.12
    add_textbox(slide, label, x + 8, y + 20, w - 16, 28, size=12, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)
    return step


def add_flow_arrow(slide, x, y, w, h):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, p(x), p(y), p(w), p(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLORS["accent"]
    arrow.fill.transparency = 0.10
    arrow.line.fill.background()
    return arrow


def build_presentation(out_path: Path, dashboard_img: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    sw = prs.slide_width
    sh = prs.slide_height

    # Slide 1
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title_block(s, "AI Smart Inventory Management System", "Demand Forecasting using Machine Learning")
    add_chip(s, "Team Members: Name 1 | Name 2 | Name 3 | Name 4", 44, 164, 520, 38)
    add_card(s, 44, 228, 420, 220, "Project Scope", "End-to-end ML workflow for inventory demand forecasting, deployment, and MLOps.")
    add_card(s, 492, 228, 424, 220, "Technology Stack", "Python, Random Forest, XGBoost, FastAPI, Streamlit, Docker, GitHub Actions")

    # Slide 2
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title_block(s, "Problem Statement")
    add_card(s, 44, 146, 872, 244, "Key Challenges", "Overstocking -> high cost\nUnderstocking -> lost sales\nNo accurate demand prediction")
    add_chip(s, "Retail needs intelligent demand forecasting", 44, 418, 520, 38, fill_key="accent")

    # Slide 3
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title_block(s, "Solution")
    add_card(s, 44, 156, 274, 250, "ML-based Prediction", "Forecast demand with supervised learning to support better stock decisions.")
    add_card(s, 343, 156, 274, 250, "Feature-rich Inputs", "Uses historical sales and engineered features to improve predictive power.")
    add_card(s, 642, 156, 274, 250, "Automated E2E Pipeline", "Integrated workflow from training to deployment and monitoring.")

    # Slide 4
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title_block(s, "Dataset", "Time-series retail data")
    add_card(s, 44, 160, 200, 200, "Store", "Store-level sales context")
    add_card(s, 264, 160, 200, 200, "Product Family", "Category-level behavior")
    add_card(s, 484, 160, 200, 200, "Sales", "Historical target signal")
    add_card(s, 704, 160, 212, 200, "Promotion", "Demand uplift driver")
    add_chip(s, "Dataset versioned using DVC", 44, 400, 330, 38)

    # Slide 5
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title_block(s, "Preprocessing")
    add_card(s, 44, 160, 280, 220, "Data Quality", "Missing values handled")
    add_card(s, 340, 160, 280, 220, "Time Consistency", "Sorted by time")
    add_card(s, 636, 160, 280, 220, "Encoding", "Encoding done")
    add_chip(s, "Same preprocessing used in training & inference", 44, 410, 520, 38, fill_key="accent2")

    # Slide 6
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title_block(s, "Feature Engineering")
    add_card(s, 44, 150, 210, 216, "Lag Features", "Lag features -> past sales")
    add_card(s, 272, 150, 210, 216, "Rolling Mean", "Rolling mean -> trend")
    add_card(s, 500, 150, 210, 216, "Rolling Std", "Rolling std -> volatility")
    add_card(s, 728, 150, 188, 216, "Time Features", "Time features -> seasonality")
    add_chip(s, "Core strength of the model", 44, 396, 310, 38, fill_key="accent")

    # Slide 7
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title_block(s, "Models Used")
    add_card(s, 120, 176, 320, 220, "Random Forest", "Ensemble baseline model")
    add_card(s, 520, 176, 320, 220, "XGBoost", "Gradient boosted decision trees")
    add_chip(s, "Compared using RMSE", 360, 420, 240, 38)

    # Slide 8
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title_block(s, "Results")
    add_chip(s, "XGBoost -> best performance", 44, 126, 360, 34, fill_key="accent2")
    add_chip(s, "Lower RMSE", 420, 126, 180, 34)
    add_chip(s, "Handles non-linear patterns better", 616, 126, 300, 34)
    chart_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p(120), p(190), p(720), p(250))
    chart_box.fill.solid()
    chart_box.fill.fore_color.rgb = COLORS["panel"]
    chart_box.line.color.rgb = COLORS["line"]
    rf = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p(190), p(275), p(460), p(38))
    rf.fill.solid()
    rf.fill.fore_color.rgb = COLORS["bad"]
    rf.line.fill.background()
    add_textbox(s, "Random Forest RMSE (higher)", 194, 248, 260, 20, size=11, color=COLORS["muted"])
    xg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p(190), p(340), p(350), p(38))
    xg.fill.solid()
    xg.fill.fore_color.rgb = COLORS["accent2"]
    xg.line.fill.background()
    add_textbox(s, "XGBoost RMSE (lower)", 194, 313, 240, 20, size=11, color=COLORS["muted"])

    # Slide 9
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title_block(s, "Pipeline Architecture")
    start_x = 28
    y = 230
    step_w = 150
    step_h = 80
    gap = 36
    aw = 26
    add_flow_step(s, start_x, y, step_w, step_h, "Data")
    add_flow_arrow(s, start_x + step_w + 5, y + 26, aw, 28)
    add_flow_step(s, start_x + (step_w + gap) * 1, y, step_w, step_h, "Preprocess")
    add_flow_arrow(s, start_x + (step_w + gap) * 1 + step_w + 5, y + 26, aw, 28)
    add_flow_step(s, start_x + (step_w + gap) * 2, y, step_w, step_h, "Features")
    add_flow_arrow(s, start_x + (step_w + gap) * 2 + step_w + 5, y + 26, aw, 28)
    add_flow_step(s, start_x + (step_w + gap) * 3, y, step_w, step_h, "Train")
    add_flow_arrow(s, start_x + (step_w + gap) * 3 + step_w + 5, y + 26, aw, 28)
    add_flow_step(s, start_x + (step_w + gap) * 4, y, step_w, step_h, "Save Model", accent=True)
    add_chip(s, "Fully automated using pipeline.py", 44, 408, 360, 38, fill_key="accent")

    # Slide 10
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title_block(s, "Inference Pipeline")
    start_x = 80
    y = 220
    step_w = 170
    step_h = 90
    gap = 48
    aw = 32
    add_flow_step(s, start_x, y, step_w, step_h, "Input")
    add_flow_arrow(s, start_x + step_w + 8, y + 30, aw, 30)
    add_flow_step(s, start_x + (step_w + gap) * 1, y, step_w, step_h, "Preprocess")
    add_flow_arrow(s, start_x + (step_w + gap) * 1 + step_w + 8, y + 30, aw, 30)
    add_flow_step(s, start_x + (step_w + gap) * 2, y, step_w, step_h, "Features")
    add_flow_arrow(s, start_x + (step_w + gap) * 2 + step_w + 8, y + 30, aw, 30)
    add_flow_step(s, start_x + (step_w + gap) * 3, y, step_w, step_h, "Predict", accent=True)
    add_chip(s, "Same logic reused -> no data leakage", 44, 408, 410, 38, fill_key="accent2")

    # Slide 11
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title_block(s, "FastAPI Deployment")
    add_card(
        s,
        44,
        164,
        430,
        242,
        "API Contract",
        "Endpoint: /predict\nJSON input\nReturns prediction\nSchema validation using Pydantic",
    )
    json_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p(500), p(164), p(416), p(242))
    json_card.fill.solid()
    json_card.fill.fore_color.rgb = RGBColor(13, 22, 35)
    json_card.line.color.rgb = COLORS["line"]
    json_card.adjustments[0] = 0.07
    add_textbox(
        s,
        '{\n  "store": 12,\n  "family": "Dairy",\n  "promotion": 1\n}\n\n-> { "prediction": 1248.6 }',
        518,
        184,
        380,
        210,
        size=13,
        color=COLORS["accent"],
    )

    # Slide 12
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title_block(s, "MLOps")
    add_card(s, 44, 164, 280, 236, "Docker", "Containerization")
    add_card(s, 340, 164, 280, 236, "Logging", "Tracking predictions")
    add_card(s, 636, 164, 280, 236, "CI/CD", "Automation with GitHub Actions")
    add_chip(s, "This slide gives extra marks", 44, 420, 290, 38, fill_key="warn")

    # Slide 13
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title_block(s, "Dashboard")
    add_chip(s, "Streamlit SaaS UI | KPI cards | Graph + prediction", 44, 118, 530, 34)
    frame = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p(44), p(160), p(872), p(334))
    frame.fill.solid()
    frame.fill.fore_color.rgb = COLORS["panel"]
    frame.line.color.rgb = COLORS["line"]
    frame.adjustments[0] = 0.03
    if dashboard_img.exists():
        s.shapes.add_picture(str(dashboard_img), p(58), p(174), p(844), p(306))
    else:
        add_textbox(s, "Screenshot here", 390, 305, 200, 30, size=18, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)

    # Slide 14
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_textbox(s, "Live Demo", 0, 210, 960, 80, size=54, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)
    add_textbox(s, "(Switch to app here)", 0, 286, 960, 30, size=16, color=COLORS["muted"], align=PP_ALIGN.CENTER)

    # Slide 15
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title_block(s, "Conclusion")
    add_card(s, 84, 180, 250, 220, "End-to-end ML system", "Integrated workflow from data to prediction")
    add_card(s, 354, 180, 250, 220, "Production-ready", "Deployment and MLOps are included")
    add_card(s, 624, 180, 250, 220, "Scalable", "Designed for extension and reuse")

    # Slide 16
    s = prs.slides.add_slide(blank)
    add_background(s, sw, sh)
    add_title_block(s, "Future Work")
    add_card(s, 80, 180, 250, 220, "External Signals", "Add weather/holiday data")
    add_card(s, 354, 180, 250, 220, "Advanced Models", "Use deep learning (LSTM)")
    add_card(s, 628, 180, 250, 220, "Real-time Integration", "Real-time integration")

    prs.save(str(out_path))


if __name__ == "__main__":
    workspace = Path.cwd()
    output = workspace / "AI_Smart_Inventory_Management_System.pptx"
    dashboard = workspace / "dashboard_mock.png"
    build_presentation(output, dashboard)
    print(f"Created: {output}")
